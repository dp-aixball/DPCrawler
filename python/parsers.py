"""
Document parsers for DPCrawler
Handles conversion from HTML and binary formats to Markdown.
"""
import io
import re
from bs4 import BeautifulSoup

try:
    import html2text
    HTML2TEXT_AVAILABLE = True
except ImportError:
    HTML2TEXT_AVAILABLE = False

try:
    import trafilatura
    TRAFILATURA_AVAILABLE = True
except ImportError:
    TRAFILATURA_AVAILABLE = False

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import docx as python_docx
except ImportError:
    python_docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    import olefile
except ImportError:
    olefile = None

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None


def extract_body_content(text: str) -> str:
    """Extract main content: keep '当前位置：' line and everything after, remove footer"""
    lines = text.split('\n')

    # Find '当前位置：' marker
    start_idx = -1
    for i, line in enumerate(lines):
        if '当前位置：' in line.strip() or '当前位置:' in line.strip():
            start_idx = i
            break

    # If found, keep from that line onwards
    if start_idx >= 0:
        lines = lines[start_idx:]
        # Merge breadcrumb lines into first line: '当前位置：首页 > 栏目'
        breadcrumb = lines[0].strip()
        merge_end = 1
        while merge_end < len(lines):
            stripped = lines[merge_end].strip()
            # compute logical text length ignoring URL parts of markdown links
            text_only = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', stripped)
            if stripped == '' or text_only == '>' or text_only == '\>' or (len(text_only) < 30 and not any(c in text_only for c in '。，；')):
                if stripped:
                    breadcrumb += ' ' + stripped
                merge_end += 1
            else:
                break
        lines = [breadcrumb] + lines[merge_end:]

    # Remove footer: find common footer markers
    footer_markers = ['关于我们', '版权所有', '京ICP备', '京公网安备', '区招生考试中心']
    end_idx = len(lines)
    for i, line in enumerate(lines):
        if i == 0:
            continue  # Skip '当前位置' line itself
        stripped = line.strip()
        for marker in footer_markers:
            if marker in stripped:
                end_idx = i
                break
        if end_idx != len(lines):
            break
    lines = lines[:end_idx]

    # Clean up: remove excessive blank lines
    result = '\n'.join(lines).strip()
    result = re.sub(r'\n{3,}', '\n\n', result)

    return result

def _table_to_markdown(table) -> str:
    """Convert a BeautifulSoup <table> element directly to Markdown table format.
    This is more reliable than html2text for complex/styled tables.
    Uses a 2D grid to correctly handle rowspan and colspan alignment."""
    grid = {}
    max_col = -1
    max_row = -1
    
    trs = table.find_all('tr')
    for r_idx, tr in enumerate(trs):
        c_idx = 0
        for cell in tr.find_all(['td', 'th']):
            # Find next available column slot in this row
            while grid.get((r_idx, c_idx)) is not None:
                c_idx += 1
            
            # Preserve line breaks for markdown cells
            for br in cell.find_all(['br', 'hr']):
                br.replace_with(' <br> ')
            for p_div in cell.find_all(['p', 'div']):
                p_div.insert_after(' <br> ')
                p_div.unwrap()
            
            # Get pure text, collapse whitespace
            text = cell.get_text(strip=True)
            text = re.sub(r'\s+', ' ', text)
            
            # Clean up redundant <br> tags
            text = re.sub(r'(<br>\s*)+', '<br>', text)
            text = text.strip('<br>').strip()
            
            # Escape pipe characters inside cell text
            text = text.replace('|', '\\|')
            
            # Handle cell spanning
            try:
                colspan = int(cell.get('colspan', 1))
            except (ValueError, TypeError):
                colspan = 1
            try:
                rowspan = int(cell.get('rowspan', 1))
            except (ValueError, TypeError):
                rowspan = 1
                
            # Fill grid with cell content and empty spans
            for r in range(r_idx, r_idx + rowspan):
                for c in range(c_idx, c_idx + colspan):
                    grid[(r, c)] = text if (r == r_idx and c == c_idx) else ""
                    max_col = max(max_col, c)
                    max_row = max(max_row, r)
            
            c_idx += colspan

    if max_row < 0 or max_col < 0:
        return ""

    rows = []
    for r in range(max_row + 1):
        row = []
        for c in range(max_col + 1):
            row.append(grid.get((r, c), ""))
        rows.append(row)

    if not rows:
        return ''

    # Normalize column count (pad short rows)
    max_cols = max(len(r) for r in rows)
    for row in rows:
        while len(row) < max_cols:
            row.append('')

    # Build markdown table
    lines = []
    # First row as header
    lines.append('| ' + ' | '.join(rows[0]) + ' |')
    # Separator
    lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
    # Data rows
    for row in rows[1:]:
        lines.append('| ' + ' | '.join(row) + ' |')

    return '\n'.join(lines)


def _convert_tables_to_markdown(html: str) -> tuple:
    """Replace all <table> elements in HTML with their Markdown equivalents.
    This ensures tables are correctly formatted regardless of html2text behavior."""
    soup = BeautifulSoup(html, "html.parser")
    placeholders = {}
    for i, table in enumerate(soup.find_all('table')):
        md_table = _table_to_markdown(table)
        if md_table:
            # Replace the table with a text placeholder to be swapped after html2text
            placeholder_key = f"___TABLE_PLACEHOLDER_{i}___"
            placeholders[placeholder_key] = md_table
            
            marker = soup.new_string(f"\n\n{placeholder_key}\n\n")
            table.replace_with(marker)
        else:
            table.decompose()
    return str(soup), placeholders


def _clean_html_inline_tags(html: str) -> str:
    """Clean HTML at DOM level: unwrap meaningless inline tags.
    The HTML block structure (<p>, <div>, etc.) already defines the correct
    paragraph/line breaks - we just need to remove inline noise tags.
    """
    soup = BeautifulSoup(html, "html.parser")

    # Unwrap inline style tags that carry no semantic meaning for RAG
    for tag_name in ['span', 'font', 'u', 'b', 'strong', 'i', 'em']:
        for tag in soup.find_all(tag_name):
            tag.unwrap()

    return str(soup)


def html_to_markdown(html: str, base_url: str = "") -> str:
    """Convert HTML to Markdown format"""
    # Clean up complex table cells before conversion
    html, table_placeholders = _convert_tables_to_markdown(html)
    # Remove meaningless inline tags that cause unwanted line breaks
    html = _clean_html_inline_tags(html)
    
    result_md = ""
    if HTML2TEXT_AVAILABLE:
        h = html2text.HTML2Text()
        h.baseurl = base_url
        h.ignore_links = False
        h.ignore_images = True   # Removed images for RAG optimization
        h.body_width = 0  # Don't wrap lines
        h.bypass_tables = False
        h.pad_tables = True  # Pad table cells for alignment
        result_md = h.handle(html)
    else:
        # Fallback: simple HTML tag removal
        soup = BeautifulSoup(html, "html.parser")
        result_md = soup.get_text(separator="\n", strip=True)
        
    # Restore tables from placeholders
    for placeholder, md_table in table_placeholders.items():
        result_md = result_md.replace(placeholder, "\n\n" + md_table + "\n\n")
        
    return result_md


def extract_true_title(html: str, default_title: str) -> str:
    """Heuristic fallback to extract true document title from badly coded sites."""
    try:
        soup = BeautifulSoup(html, "html.parser")
        
        # 1. Classic article title identifiers
        identifiers = [
            'arti_title', 'title', 'news_title', 'article-title', 'con_title', 
            'bt', 'biaoti', 'tit', 'article_title', 'arti-name', 'titleTxt',
            'info-ctit', 'info_title', 'articleTitle', 'NewsTitle'
        ]
        for ident in identifiers:
            t_elem = soup.find(class_=ident) or soup.find(id=ident)
            if t_elem:
                text = t_elem.get_text(separator=' ', strip=True)
                if len(text) > 4:
                    return text

        # 2. Header tags
        for tag in ['h1', 'h2']:
            for elem in soup.find_all(tag):
                text = elem.get_text(separator=' ', strip=True)
                if 6 < len(text) < 100:
                    if not default_title or text not in default_title:
                        return text

        # 3. Readability short title
        try:
            from readability import Document
            doc = Document(html)
            r_title = doc.short_title().strip()
            if len(r_title) > 4 and r_title != default_title:
                return r_title
        except:
            pass
    except Exception:
        pass
    return default_title


def universal_html_extract(html: str, base_url: str = "") -> dict:
    """
    Extract main content and metadata using Trafilatura for universal support.
    Returns a dict: {"content": str, "metadata": dict}
    Will fallback to custom html_to_markdown if Trafilatura fails or returns too little.
    """
    fallback_result = None

    def _do_fallback():
        content = html_to_markdown(html, base_url)
        content = extract_body_content(content)
        return {"content": content, "metadata": {}}

    if not TRAFILATURA_AVAILABLE:
        return _do_fallback()

    try:
        # Pre-convert complex tables manually to preserve rowspan/colspan,
        # then let trafilatura handle the rest. Trafilatura is good but our table parser handles edge-cases.
        html_processed, table_placeholders = _convert_tables_to_markdown(html)
        
        md_content = trafilatura.extract(html_processed, url=base_url, output_format="markdown", include_tables=True, include_images=False, include_comments=False)
        metadata = trafilatura.extract_metadata(html_processed, default_url=base_url)
        
        if md_content and len(md_content.strip()) > 50:
            # Restore tables
            for placeholder, md_table in table_placeholders.items():
                md_content = md_content.replace(placeholder, "\n\n" + md_table + "\n\n")

            meta_dict = {}
            if hasattr(metadata, 'as_dict'):
                meta_dict = metadata.as_dict()
            elif isinstance(metadata, dict):
                meta_dict = metadata
            
            # Clean up metadata (remove Nones)
            clean_meta = {k: v for k, v in meta_dict.items() if v}
            
            title_text = clean_meta.get('title', '').strip()
            
            # Cross-verify and fix title using heuristics
            better_title = extract_true_title(html_processed, title_text)
            if better_title:
                clean_meta['title'] = better_title
                title_text = better_title
            
            # Re-inject title into body if not already present
            if title_text:
                md_head = md_content[:200].replace('#', '').strip()
                if title_text not in md_head and not md_head.startswith(title_text[:10]):
                    md_content = f"# {title_text}\n\n{md_content}"
            
            return {"content": md_content, "metadata": clean_meta}
        else:
            return _do_fallback()
    except Exception as e:
        print(f"  -> Trafilatura parsing error: {e}")
        return _do_fallback()


def doc_to_markdown(data: bytes) -> str:
    """Convert old .doc (OLE binary) to Markdown using olefile (piece table extraction)"""
    if olefile is None:
        return "[.doc conversion unavailable: install olefile]"
    
    try:
        ole = olefile.OleFileIO(io.BytesIO(data))
        
        # Check if it's a valid Word document
        if not ole.exists('WordDocument'):
            return "[.doc conversion failed: not a valid Word document]"
        
        # Read WordDocument stream
        word_doc = ole.openstream('WordDocument').read()
        
        if len(word_doc) < 34:
            return "[.doc conversion failed: file too small]"
        
        # Check magic number
        w_ident = word_doc[0] | (word_doc[1] << 8)
        if w_ident != 0xA5EC:
            return "[.doc conversion failed: invalid Word format]"
        
        # Determine which Table stream to use (0Table or 1Table)
        flags = word_doc[0x0A] | (word_doc[0x0B] << 8)
        f_which_tbl = (flags >> 9) & 1
        table_name = '1Table' if f_which_tbl == 1 else '0Table'
        
        # Parse FIB to find CLX (piece table)
        csw = word_doc[32] | (word_doc[33] << 8)
        rg_w_end = 34 + csw * 2
        
        if len(word_doc) < rg_w_end + 2:
            return "[.doc conversion failed: invalid FIB structure]"
        
        cslw = word_doc[rg_w_end] | (word_doc[rg_w_end + 1] << 8)
        rg_lw_start = rg_w_end + 2
        
        # ccpText is at index 3 in fibRgLw
        if cslw < 4:
            return "[.doc conversion failed: no text content]"
        ccp_text = (word_doc[rg_lw_start + 12] | (word_doc[rg_lw_start + 13] << 8) | 
                    (word_doc[rg_lw_start + 14] << 16) | (word_doc[rg_lw_start + 15] << 24))
        if ccp_text == 0:
            return ""
        
        # Find fcClx/lcbClx (pair 33 in fibRgFcLcb)
        rg_lw_end = rg_lw_start + cslw * 4
        if len(word_doc) < rg_lw_end + 2:
            return "[.doc conversion failed: truncated FIB]"
        cb_rg_fc_lcb = word_doc[rg_lw_end] | (word_doc[rg_lw_end + 1] << 8)
        
        if cb_rg_fc_lcb <= 33:
            return "[.doc conversion failed: no CLX found]"
        
        fc_clx_offset = rg_lw_end + 2 + 33 * 8
        if len(word_doc) < fc_clx_offset + 8:
            return "[.doc conversion failed: truncated CLX info]"
        
        fc_clx = (word_doc[fc_clx_offset] | (word_doc[fc_clx_offset + 1] << 8) | 
                  (word_doc[fc_clx_offset + 2] << 16) | (word_doc[fc_clx_offset + 3] << 24))
        lcb_clx = (word_doc[fc_clx_offset + 4] | (word_doc[fc_clx_offset + 5] << 8) | 
                   (word_doc[fc_clx_offset + 6] << 16) | (word_doc[fc_clx_offset + 7] << 24))
        
        if lcb_clx == 0:
            return "[.doc conversion failed: empty CLX]"
        
        # Read Table stream
        if not ole.exists(table_name):
            return "[.doc conversion failed: table stream missing]"
        table_data = ole.openstream(table_name).read()
        
        if fc_clx + lcb_clx > len(table_data):
            return "[.doc conversion failed: CLX out of bounds]"
        
        # Parse CLX to find piece table
        clx = table_data[fc_clx:fc_clx + lcb_clx]
        pos = 0
        
        # Skip Prc records (0x01), find Pcdt (0x02)
        while pos < len(clx):
            if clx[pos] == 0x01:
                if pos + 3 > len(clx):
                    break
                cb = clx[pos + 1] | (clx[pos + 2] << 8)
                pos += 3 + cb
            elif clx[pos] == 0x02:
                pos += 1
                break
            else:
                break
        
        if pos + 4 > len(clx):
            return "[.doc conversion failed: no piece table]"
        
        lcb_plc_pcd = (clx[pos] | (clx[pos + 1] << 8) | 
                       (clx[pos + 2] << 16) | (clx[pos + 3] << 24))
        pos += 4
        
        if pos + lcb_plc_pcd > len(clx) or lcb_plc_pcd < 16:
            return "[.doc conversion failed: invalid piece table]"
        
        plc_pcd = clx[pos:pos + lcb_plc_pcd]
        
        # Parse pieces
        n = (lcb_plc_pcd - 4) // 12
        if n == 0:
            return "[.doc conversion failed: no pieces]"
        
        pcd_start = (n + 1) * 4
        text_parts = []
        
        for i in range(n):
            cp_start = (plc_pcd[i * 4] | (plc_pcd[i * 4 + 1] << 8) | 
                        (plc_pcd[i * 4 + 2] << 16) | (plc_pcd[i * 4 + 3] << 24))
            cp_end = (plc_pcd[(i + 1) * 4] | (plc_pcd[(i + 1) * 4 + 1] << 8) | 
                      (plc_pcd[(i + 1) * 4 + 2] << 16) | (plc_pcd[(i + 1) * 4 + 3] << 24))
            char_count = max(0, cp_end - cp_start)
            
            if char_count == 0 or cp_start >= ccp_text:
                continue
            
            effective_count = min(char_count, ccp_text - cp_start)
            
            pcd_offset = pcd_start + i * 8
            if pcd_offset + 6 > len(plc_pcd):
                break
            
            fc_raw = (plc_pcd[pcd_offset + 2] | (plc_pcd[pcd_offset + 3] << 8) | 
                      (plc_pcd[pcd_offset + 4] << 16) | (plc_pcd[pcd_offset + 5] << 24))
            is_compressed = (fc_raw & 0x40000000) != 0
            fc = fc_raw & 0x3FFFFFFF
            
            if is_compressed:
                # ANSI: 1 byte = 1 char
                byte_start = fc // 2
                byte_end = byte_start + effective_count
                if byte_end > len(word_doc):
                    continue
                
                chars = []
                for b in word_doc[byte_start:byte_end]:
                    if b == 0x0D or b == 0x0C:
                        chars.append('\n')
                    elif b == 0x07:
                        chars.append('\t')
                    elif b >= 0x20:
                        chars.append(chr(b))
                text_parts.append(''.join(chars))
            else:
                # Unicode UTF-16LE: 2 bytes = 1 char
                byte_start = fc
                byte_end = byte_start + effective_count * 2
                if byte_end > len(word_doc):
                    continue
                
                u16s = []
                for j in range(byte_start, byte_end, 2):
                    if j + 1 < len(word_doc):
                        u16s.append(word_doc[j] | (word_doc[j + 1] << 8))
                
                try:
                    decoded = ''.join(chr(c) for c in u16s)
                except:
                    decoded = ''.join(chr(c) if c < 0x10000 else '?' for c in u16s)
                
                # Normalize control chars
                cleaned = []
                for ch in decoded:
                    if ch == '\r' or ch == '\x0C':
                        cleaned.append('\n')
                    elif ch == '\x07':
                        cleaned.append('\t')
                    elif ch >= ' ' or ch == '\n' or ch == '\t':
                        cleaned.append(ch)
                text_parts.append(''.join(cleaned))
        
        ole.close()
        
        result = ''.join(text_parts).strip()
        if not result:
            return "[.doc conversion failed: no text extracted]"
        
        # Clean up excessive blank lines
        result = re.sub(r'\n{3,}', '\n\n', result)
        return result
        
    except Exception as e:
        return f"[.doc conversion error: {e}]"


def pdf_to_markdown(data: bytes) -> str:
    """Convert PDF binary data to Markdown text"""
    if pdfplumber is None:
        return "[PDF conversion unavailable: install pdfplumber]"
    lines = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages):
            # Extract tables
            tables = page.extract_tables()
            if tables:
                for table in tables:
                    if not table:
                        continue
                    # Build MD table
                    max_cols = max(len(row) for row in table if row)
                    header = table[0] if table else []
                    header = [(c or '').strip().replace('\n', ' ').replace('|', '\\|') for c in header]
                    while len(header) < max_cols:
                        header.append('')
                    lines.append('| ' + ' | '.join(header) + ' |')
                    lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
                    for row in table[1:]:
                        if not row:
                            continue
                        cells = [(c or '').strip().replace('\n', ' ').replace('|', '\\|') for c in row]
                        while len(cells) < max_cols:
                            cells.append('')
                        lines.append('| ' + ' | '.join(cells) + ' |')
                    lines.append('')
            else:
                text = page.extract_text()
                if text:
                    lines.append(text)
                    lines.append('')
    return '\n'.join(lines).strip()


def docx_to_markdown(data: bytes) -> str:
    """Convert DOCX binary data to Markdown text"""
    if python_docx is None:
        return "[DOCX conversion unavailable: install python-docx]"
    doc = python_docx.Document(io.BytesIO(data))
    lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            lines.append('')
            continue
        style = para.style.name.lower() if para.style else ''
        if 'heading 1' in style:
            lines.append(f'# {text}')
        elif 'heading 2' in style:
            lines.append(f'## {text}')
        elif 'heading 3' in style:
            lines.append(f'### {text}')
        elif 'heading 4' in style:
            lines.append(f'#### {text}')
        elif 'list' in style:
            lines.append(f'- {text}')
        else:
            lines.append(text)
        lines.append('')
    # Convert tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip().replace('\n', ' ').replace('|', '\\|') for cell in row.cells]
            rows.append(cells)
        if rows:
            max_cols = max(len(r) for r in rows)
            for r in rows:
                while len(r) < max_cols:
                    r.append('')
            lines.append('| ' + ' | '.join(rows[0]) + ' |')
            lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
            for r in rows[1:]:
                lines.append('| ' + ' | '.join(r) + ' |')
            lines.append('')
    return '\n'.join(lines).strip()


def xls_to_markdown(data: bytes) -> str:
    """Convert old .xls (OLE binary) to Markdown tables using xlrd"""
    if xlrd is None:
        return "[XLS conversion unavailable: install xlrd]"
    
    try:
        wb = xlrd.open_workbook(file_contents=data)
        lines = []
        
        for sheet_idx in range(wb.nsheets):
            sheet = wb.sheet_by_index(sheet_idx)
            lines.append(f'## {sheet.name}')
            lines.append('')
            
            if sheet.nrows == 0:
                lines.append('*空工作表*')
                lines.append('')
                continue
            
            # Build Markdown table
            max_cols = sheet.ncols
            rows = []
            for row_idx in range(sheet.nrows):
                cells = []
                for col_idx in range(max_cols):
                    cell = sheet.cell(row_idx, col_idx)
                    if cell.ctype == xlrd.XL_CELL_EMPTY:
                        cells.append('')
                    elif cell.ctype == xlrd.XL_CELL_NUMBER:
                        # Format number: remove trailing .0 for integers
                        val = cell.value
                        if val == int(val):
                            cells.append(str(int(val)))
                        else:
                            cells.append(str(val))
                    else:
                        cells.append(str(cell.value).strip().replace('|', '\\|'))
                rows.append(cells)
            
            if rows:
                # First row as header
                lines.append('| ' + ' | '.join(rows[0]) + ' |')
                lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
                for row in rows[1:]:
                    lines.append('| ' + ' | '.join(row) + ' |')
            lines.append('')
        
        return '\n'.join(lines).strip()
    except Exception as e:
        return f"[XLS conversion error: {e}]"


def xlsx_to_markdown(data: bytes) -> str:
    """Convert XLSX binary data to Markdown tables"""
    if openpyxl is None:
        return "[XLSX conversion unavailable: install openpyxl]"
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f'## {sheet.title}')
        lines.append('')
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip().replace('\n', ' ').replace('|', '\\|') if c is not None else '' for c in row]
            # Skip fully empty rows
            if not any(cells):
                continue
            # Skip merged title rows (only 1-2 cells have content, rest are empty)
            non_empty = sum(1 for c in cells if c)
            if non_empty <= 2 and len(cells) > 4:
                # Treat as a title line, not a table row
                title_text = ' '.join(c for c in cells if c)
                if title_text:
                    lines.append(f'**{title_text}**')
                    lines.append('')
                continue
            rows.append(cells)
        if rows:
            max_cols = max(len(r) for r in rows)
            for r in rows:
                while len(r) < max_cols:
                    r.append('')
            lines.append('| ' + ' | '.join(rows[0]) + ' |')
            lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
            for r in rows[1:]:
                lines.append('| ' + ' | '.join(r) + ' |')
        lines.append('')
    wb.close()
    return '\n'.join(lines).strip()


def pptx_to_markdown(data: bytes) -> str:
    """Convert PPTX binary data to Markdown text"""
    if PptxPresentation is None:
        return "[PPTX conversion unavailable: install python-pptx]"
    prs = PptxPresentation(io.BytesIO(data))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f'## Slide {i}')
        lines.append('')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip().replace('\n', ' ').replace('|', '\\|') for cell in row.cells]
                    rows.append(cells)
                if rows:
                    max_cols = max(len(r) for r in rows)
                    for r in rows:
                        while len(r) < max_cols:
                            r.append('')
                    lines.append('| ' + ' | '.join(rows[0]) + ' |')
                    lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
                    for r in rows[1:]:
                        lines.append('| ' + ' | '.join(r) + ' |')
        lines.append('')
    return '\n'.join(lines).strip()

def extract_title(soup: BeautifulSoup) -> str:
    """Extract title from HTML page"""
    # Try h1 first
    h1 = soup.find("h1")
    if h1:
        return h1.get_text(strip=True)

    # Try title tag
    title = soup.find("title")
    if title:
        return title.get_text(strip=True)

    return "Untitled"

def extract_main_html(raw_html: str, title: str = "") -> str:
    """Extract semantic main content HTML from raw webpage DOM"""
    # 优先尝试使用专业的 Readability 引擎剥离网页模板噪声（如面包屑、边栏、版权页）
    try:
        from readability import Document
        doc = Document(raw_html)
        clean_html = doc.summary(html_partial=True)
        if clean_html and len(clean_html) > 100:
            date_str = ""
            
            # 引入项目自带的老式网页大标题扫描器（专门针对类 info-ctit 这种政务网站硬编码class）
            better_title = extract_true_title(raw_html, title)
            if better_title and len(better_title) > 3 and better_title != title:
                title = better_title
                
            try:
                import trafilatura
                meta = trafilatura.extract_metadata(raw_html)
                if meta:
                    if getattr(meta, 'date', None):
                        date_str = meta.date
                    # Trafilatura 也是一重屏障
                    tmp_title = getattr(meta, 'title', None)
                    if tmp_title and len(tmp_title) > 3 and title == better_title:
                        pass # 已被 extract_true_title 精准捕获则放弃 trafilatura
                    elif tmp_title and len(tmp_title) > 3:
                        title = tmp_title
            except Exception:
                pass
                
            header_html = f"<div class=\"rag-header\">\n  <h1 class=\"rag-title\">{title}</h1>\n"
            if date_str:
                header_html += f"  <div class=\"rag-meta\"><span>发布时间：{date_str}</span></div>\n"
            header_html += "</div>\n"
            
            css = """
body { font-family: -apple-system, "PingFang SC", "Microsoft YaHei", "Segoe UI", Roboto, Helvetica, Arial, sans-serif; line-height: 1.8; max-width: 900px; margin: 40px auto; padding: 0 20px; color: #2c3e50; font-size: 16px; background-color: #fafafa; }
img { max-width: 100%; height: auto; border-radius: 6px; margin: 20px 0; box-shadow: 0 2px 8px rgba(0,0,0,0.05); }
table { border-collapse: collapse; width: 100%; margin-bottom: 24px; background: #fff; box-shadow: 0 1px 3px rgba(0,0,0,0.05); }
th, td { border: 1px solid #e2e8f0; padding: 12px 16px; text-align: left; }
th { background-color: #f8fafc; font-weight: 600; color: #475569; }
.rag-header { margin-bottom: 30px; text-align: center; border-bottom: 2px solid #edf2f7; padding-bottom: 24px; }
.rag-title { font-size: 2em; font-weight: 700; color: #1e293b; margin-bottom: 12px; line-height: 1.4; }
.rag-meta { font-size: 0.95em; color: #94a3b8; display: flex; justify-content: center; gap: 16px; }
p { margin-bottom: 1.5em; text-indent: 2em; text-align: justify; }
a { color: #3b82f6; text-decoration: none; }
a:hover { text-decoration: underline; }
"""
            return f"<!DOCTYPE html>\n<html>\n<head>\n<meta charset=\"utf-8\">\n<style>{css}</style>\n<title>{title}</title>\n</head>\n<body>\n{header_html}{clean_html}\n</body>\n</html>"
    except Exception:
        pass

    soup = BeautifulSoup(raw_html, "html.parser")
    
    # Remove unwanted tags universally
    for tag in soup.find_all(['script', 'style', 'nav', 'footer', 'header', 'aside', 'noscript', 'iframe']):
        tag.decompose()
        
    candidates = soup.find_all(['article', 'main'])
    
    if not candidates:
        for div in soup.find_all('div'):
            classes = " ".join(div.get('class', [])).lower()
            idx = div.get('id', '').lower()
            if re.search(r'content|article|main|body|txt|zoom', classes) or \
               re.search(r'content|article|main|body|txt|zoom', idx):
                candidates.append(div)
                
    if not candidates:
        divs = soup.find_all('div')
        if divs:
            candidates.append(max(divs, key=lambda d: len(d.get_text(strip=True))))
            
    if candidates:
        best_node = max(candidates, key=lambda c: len(c.get_text(strip=True)))
        return f"<!DOCTYPE html>\n<html>\n<head>\n<meta charset=\"utf-8\">\n<style>\n  body {{ font-family: -apple-system, BlinkMacSystemFont, \"Segoe UI\", Roboto, Helvetica, Arial, sans-serif; line-height: 1.6; max-width: 800px; margin: 40px auto; padding: 0 20px; color: #333; }}\n  img {{ max-width: 100%; height: auto; }}\n  table {{ border-collapse: collapse; width: 100%; margin-bottom: 20px; }}\n  th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}\n</style>\n<title>{title}</title>\n</head>\n<body>\n{best_node.prettify()}\n</body>\n</html>"
    
    body = soup.find('body')
    return str(body) if body else raw_html
