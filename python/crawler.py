"""
Main crawler module for DPCrawler
"""
import os
import sys
import json
import time
import threading
from collections import deque
from concurrent.futures import ThreadPoolExecutor, as_completed
from urllib.parse import urljoin, urlparse
from typing import Set, Optional
import re

import requests
from bs4 import BeautifulSoup
import io
try:
    import html2text
    HTML2TEXT_AVAILABLE = True
except ImportError:
    HTML2TEXT_AVAILABLE = False

# Optional imports for document conversion (graceful degradation)
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    import docx as python_docx
except ImportError:
    python_docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None

from config import CrawlerConfig
from storage import StorageManager, CrawlResult
from dataclasses import asdict


class WebCrawler:
    """Web crawler for RAG knowledge collection"""

    MAX_PAGES = 10000  # Maximum pages per crawl session

    CONCURRENT_WORKERS = 8  # Parallel crawl workers

    def __init__(self, config: CrawlerConfig):
        self.config = config
        self.session = requests.Session()
        self.session.headers.update({
            "User-Agent": "Mozilla/5.0 (compatible; DPCrawler/1.0; +https://github.com/dpcrawler)"
        })
        self.visited_urls: Set[str] = set()
        self.new_files: list = []
        self.updated_files: list = []
        self.unchanged_count: int = 0
        self.error_count: int = 0
        self.page_count: int = 0
        self.storage: Optional[StorageManager] = None
        self.current_subdir: str = ""
        self.base_urls: list = []  # base URLs for scope checking
        self.bfs_queue: deque = deque()  # BFS queue (thread-safe with lock)
        self.url_depths: dict = {}  # Track depth of each URL
        self._lock = threading.Lock()  # Protect shared state

    @staticmethod
    def _normalize_url(url: str) -> str:
        """Normalize URL: remove fragment, sort query params, strip trailing slash"""
        parsed = urlparse(url)
        # Remove fragment
        path = parsed.path.rstrip('/') or '/'
        # Keep query but drop fragment
        normalized = f"{parsed.scheme}://{parsed.netloc}{path}"
        if parsed.query:
            normalized += f"?{parsed.query}"
        return normalized

    def _is_in_scope(self, url: str) -> bool:
        """Check if URL is within the scope of any base URL (same domain + path prefix)"""
        parsed = urlparse(url)
        url_path = parsed.path or '/'
        for base in self.base_urls:
            bp = urlparse(base)
            base_path = bp.path.rstrip('/') or '/'
            # Normalize: /foo should match /foo and /foo/bar but not /foobar
            if parsed.netloc == bp.netloc:
                if url_path == base_path or url_path.startswith(base_path + '/') or base_path == '/':
                    return True
        return len(self.base_urls) == 0

    def should_process_url(self, url: str) -> bool:
        """Check if URL should be processed based on config and scope"""
        # Scope check: must be under base URL path
        if not self._is_in_scope(url):
            return False

        parsed = urlparse(url)
        path = parsed.path.lower()

        # Check if URL has allowed extension
        for ext in self.config.file_extensions:
            if path.endswith(ext):
                return True

        # Also process HTML pages by default
        if not "." in path.split("/")[-1]:
            return True

        return False

    def extract_body_content(self, text: str) -> str:
        """Extract main content: keep '当前位置：' line and everything after, remove footer"""
        lines = text.split('\n')

        # Find '当前位置：' marker
        start_idx = -1
        for i, line in enumerate(lines):
            if '当前位置：' in line.strip() or '当前位置:' in line.strip():
                start_idx = i
                break

        # If found, keep from that line onwards
        if start_idx >= 0:
            lines = lines[start_idx:]
            # Merge breadcrumb lines into first line: '当前位置：首页 > 栏目'
            breadcrumb = lines[0].strip()
            merge_end = 1
            while merge_end < len(lines):
                stripped = lines[merge_end].strip()
                # compute logical text length ignoring URL parts of markdown links
                text_only = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', stripped)
                if stripped == '' or text_only == '>' or text_only == '\>' or (len(text_only) < 30 and not any(c in text_only for c in '。，；')):
                    if stripped:
                        breadcrumb += ' ' + stripped
                    merge_end += 1
                else:
                    break
            lines = [breadcrumb] + lines[merge_end:]

        # Remove footer: find common footer markers
        footer_markers = ['关于我们', '版权所有', '京ICP备', '京公网安备', '区招生考试中心']
        end_idx = len(lines)
        for i, line in enumerate(lines):
            if i == 0:
                continue  # Skip '当前位置' line itself
            stripped = line.strip()
            for marker in footer_markers:
                if marker in stripped:
                    end_idx = i
                    break
            if end_idx != len(lines):
                break
        lines = lines[:end_idx]

        # Clean up: remove excessive blank lines
        result = '\n'.join(lines).strip()
        result = re.sub(r'\n{3,}', '\n\n', result)

        return result

    def _table_to_markdown(self, table) -> str:
        """Convert a BeautifulSoup <table> element directly to Markdown table format.
        This is more reliable than html2text for complex/styled tables.
        Uses a 2D grid to correctly handle rowspan and colspan alignment."""
        grid = {}
        max_col = -1
        max_row = -1
        
        trs = table.find_all('tr')
        for r_idx, tr in enumerate(trs):
            c_idx = 0
            for cell in tr.find_all(['td', 'th']):
                # Find next available column slot in this row
                while grid.get((r_idx, c_idx)) is not None:
                    c_idx += 1
                
                # Preserve line breaks for markdown cells
                for br in cell.find_all(['br', 'hr']):
                    br.replace_with(' <br> ')
                for p_div in cell.find_all(['p', 'div']):
                    p_div.insert_after(' <br> ')
                    p_div.unwrap()
                
                # Get pure text, collapse whitespace
                text = cell.get_text(strip=True)
                text = re.sub(r'\s+', ' ', text)
                
                # Clean up redundant <br> tags
                text = re.sub(r'(<br>\s*)+', '<br>', text)
                text = text.strip('<br>').strip()
                
                # Escape pipe characters inside cell text
                text = text.replace('|', '\\|')
                
                # Handle cell spanning
                try:
                    colspan = int(cell.get('colspan', 1))
                except (ValueError, TypeError):
                    colspan = 1
                try:
                    rowspan = int(cell.get('rowspan', 1))
                except (ValueError, TypeError):
                    rowspan = 1
                    
                # Fill grid with cell content and empty spans
                for r in range(r_idx, r_idx + rowspan):
                    for c in range(c_idx, c_idx + colspan):
                        grid[(r, c)] = text if (r == r_idx and c == c_idx) else ""
                        max_col = max(max_col, c)
                        max_row = max(max_row, r)
                
                c_idx += colspan

        if max_row < 0 or max_col < 0:
            return ""

        rows = []
        for r in range(max_row + 1):
            row = []
            for c in range(max_col + 1):
                row.append(grid.get((r, c), ""))
            rows.append(row)

        if not rows:
            return ''

        # Normalize column count (pad short rows)
        max_cols = max(len(r) for r in rows)
        for row in rows:
            while len(row) < max_cols:
                row.append('')

        # Build markdown table
        lines = []
        # First row as header
        lines.append('| ' + ' | '.join(rows[0]) + ' |')
        # Separator
        lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
        # Data rows
        for row in rows[1:]:
            lines.append('| ' + ' | '.join(row) + ' |')

        return '\n'.join(lines)

    def _convert_tables_to_markdown(self, html: str) -> tuple:
        """Replace all <table> elements in HTML with their Markdown equivalents.
        This ensures tables are correctly formatted regardless of html2text behavior."""
        soup = BeautifulSoup(html, "html.parser")
        placeholders = {}
        for i, table in enumerate(soup.find_all('table')):
            md_table = self._table_to_markdown(table)
            if md_table:
                # Replace the table with a text placeholder to be swapped after html2text
                placeholder_key = f"___TABLE_PLACEHOLDER_{i}___"
                placeholders[placeholder_key] = md_table
                
                marker = soup.new_string(f"\n\n{placeholder_key}\n\n")
                table.replace_with(marker)
            else:
                table.decompose()
        return str(soup), placeholders

    def _clean_html_inline_tags(self, html: str) -> str:
        """Clean HTML at DOM level: unwrap meaningless inline tags.
        The HTML block structure (<p>, <div>, etc.) already defines the correct
        paragraph/line breaks - we just need to remove inline noise tags.
        """
        soup = BeautifulSoup(html, "html.parser")

        # Unwrap inline style tags that carry no semantic meaning for RAG
        for tag_name in ['span', 'font', 'u', 'b', 'strong', 'i', 'em']:
            for tag in soup.find_all(tag_name):
                tag.unwrap()

        # Removed the logic that flattens <p> inside <p> as it destroys line breaks.
        # html2text natively handles nested blocks very well.

        return str(soup)

    def html_to_markdown(self, html: str, base_url: str = "") -> str:
        """Convert HTML to Markdown format"""
        # Clean up complex table cells before conversion
        html, table_placeholders = self._convert_tables_to_markdown(html)
        # Remove meaningless inline tags that cause unwanted line breaks
        html = self._clean_html_inline_tags(html)
        
        result_md = ""
        if HTML2TEXT_AVAILABLE:
            h = html2text.HTML2Text()
            h.baseurl = base_url
            h.ignore_links = False
            h.ignore_images = True   # Removed images for RAG optimization
            h.body_width = 0  # Don't wrap lines
            h.bypass_tables = False
            h.pad_tables = True  # Pad table cells for alignment
            result_md = h.handle(html)
        else:
            # Fallback: simple HTML tag removal
            soup = BeautifulSoup(html, "html.parser")
            result_md = soup.get_text(separator="\n", strip=True)
            
        # Restore tables from placeholders
        for placeholder, md_table in table_placeholders.items():
            result_md = result_md.replace(placeholder, "\n\n" + md_table + "\n\n")
            
        return result_md

    # ---- Document format converters ----

    @staticmethod
    def pdf_to_markdown(data: bytes) -> str:
        """Convert PDF binary data to Markdown text"""
        if pdfplumber is None:
            return "[PDF conversion unavailable: install pdfplumber]"
        lines = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for i, page in enumerate(pdf.pages):
                # Extract tables
                tables = page.extract_tables()
                if tables:
                    for table in tables:
                        if not table:
                            continue
                        # Build MD table
                        max_cols = max(len(row) for row in table if row)
                        header = table[0] if table else []
                        header = [(c or '').strip().replace('\n', ' ').replace('|', '\\|') for c in header]
                        while len(header) < max_cols:
                            header.append('')
                        lines.append('| ' + ' | '.join(header) + ' |')
                        lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
                        for row in table[1:]:
                            if not row:
                                continue
                            cells = [(c or '').strip().replace('\n', ' ').replace('|', '\\|') for c in row]
                            while len(cells) < max_cols:
                                cells.append('')
                            lines.append('| ' + ' | '.join(cells) + ' |')
                        lines.append('')
                else:
                    text = page.extract_text()
                    if text:
                        lines.append(text)
                        lines.append('')
        return '\n'.join(lines).strip()

    @staticmethod
    def docx_to_markdown(data: bytes) -> str:
        """Convert DOCX binary data to Markdown text"""
        if python_docx is None:
            return "[DOCX conversion unavailable: install python-docx]"
        doc = python_docx.Document(io.BytesIO(data))
        lines = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                lines.append('')
                continue
            style = para.style.name.lower() if para.style else ''
            if 'heading 1' in style:
                lines.append(f'# {text}')
            elif 'heading 2' in style:
                lines.append(f'## {text}')
            elif 'heading 3' in style:
                lines.append(f'### {text}')
            elif 'heading 4' in style:
                lines.append(f'#### {text}')
            elif 'list' in style:
                lines.append(f'- {text}')
            else:
                lines.append(text)
            lines.append('')
        # Convert tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip().replace('\n', ' ').replace('|', '\\|') for cell in row.cells]
                rows.append(cells)
            if rows:
                max_cols = max(len(r) for r in rows)
                for r in rows:
                    while len(r) < max_cols:
                        r.append('')
                lines.append('| ' + ' | '.join(rows[0]) + ' |')
                lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
                for r in rows[1:]:
                    lines.append('| ' + ' | '.join(r) + ' |')
                lines.append('')
        return '\n'.join(lines).strip()

    @staticmethod
    def xlsx_to_markdown(data: bytes) -> str:
        """Convert XLSX binary data to Markdown tables"""
        if openpyxl is None:
            return "[XLSX conversion unavailable: install openpyxl]"
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f'## {sheet.title}')
            lines.append('')
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c).strip().replace('\n', ' ').replace('|', '\\|') if c is not None else '' for c in row]
                # Skip fully empty rows
                if not any(cells):
                    continue
                # Skip merged title rows (only 1-2 cells have content, rest are empty)
                non_empty = sum(1 for c in cells if c)
                if non_empty <= 2 and len(cells) > 4:
                    # Treat as a title line, not a table row
                    title_text = ' '.join(c for c in cells if c)
                    if title_text:
                        lines.append(f'**{title_text}**')
                        lines.append('')
                    continue
                rows.append(cells)
            if rows:
                max_cols = max(len(r) for r in rows)
                for r in rows:
                    while len(r) < max_cols:
                        r.append('')
                lines.append('| ' + ' | '.join(rows[0]) + ' |')
                lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
                for r in rows[1:]:
                    lines.append('| ' + ' | '.join(r) + ' |')
            lines.append('')
        wb.close()
        return '\n'.join(lines).strip()

    @staticmethod
    def pptx_to_markdown(data: bytes) -> str:
        """Convert PPTX binary data to Markdown text"""
        if PptxPresentation is None:
            return "[PPTX conversion unavailable: install python-pptx]"
        prs = PptxPresentation(io.BytesIO(data))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f'## Slide {i}')
            lines.append('')
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip().replace('\n', ' ').replace('|', '\\|') for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        max_cols = max(len(r) for r in rows)
                        for r in rows:
                            while len(r) < max_cols:
                                r.append('')
                        lines.append('| ' + ' | '.join(rows[0]) + ' |')
                        lines.append('| ' + ' | '.join(['---'] * max_cols) + ' |')
                        for r in rows[1:]:
                            lines.append('| ' + ' | '.join(r) + ' |')
            lines.append('')
        return '\n'.join(lines).strip()

    # Content type to converter mapping
    BINARY_CONVERTERS = {
        'application/pdf': ('pdf_to_markdown', '.pdf'),
        'application/msword': ('docx_to_markdown', '.doc'),
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': ('docx_to_markdown', '.docx'),
        'application/vnd.ms-excel': ('xlsx_to_markdown', '.xls'),
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': ('xlsx_to_markdown', '.xlsx'),
        'application/vnd.ms-powerpoint': ('pptx_to_markdown', '.ppt'),
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': ('pptx_to_markdown', '.pptx'),
    }

    # Text types that can be stored as-is (wrapped in MD code block or plain)
    PLAINTEXT_TYPES = {
        'text/plain', 'text/csv', 'text/xml', 'application/json',
        'application/xml', 'text/markdown', 'application/rtf',
    }

    # URL extensions for binary files
    BINARY_EXTENSIONS = {'.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'}
    PLAINTEXT_EXTENSIONS = {'.txt', '.csv', '.json', '.xml', '.yaml', '.yml', '.md', '.rst', '.tex', '.log', '.rtf'}

    def extract_title(self, soup: BeautifulSoup) -> str:
        """Extract title from HTML page"""
        # Try h1 first
        h1 = soup.find("h1")
        if h1:
            return h1.get_text(strip=True)

        # Try title tag
        title = soup.find("title")
        if title:
            return title.get_text(strip=True)

        return "Untitled"

    def extract_links(self, soup: BeautifulSoup, base_url: str) -> list:
        """Extract all links from HTML page, filtered by scope"""
        links = set()
        for a in soup.find_all("a", href=True):
            href = a["href"].strip()
            if not href or href.startswith(('javascript:', 'mailto:', 'tel:', '#')):
                continue
            full_url = urljoin(base_url, href)
            normalized = self._normalize_url(full_url)
            # Must be same domain AND within base URL scope
            if urlparse(normalized).netloc == urlparse(base_url).netloc and self._is_in_scope(normalized):
                links.add(normalized)
        return list(links)

    def process_page(self, url: str) -> bool:
        """Process a single page (called from BFS queue)"""
        normalized = self._normalize_url(url)

        # Thread-safe check and claim
        with self._lock:
            if self.page_count >= self.MAX_PAGES:
                return False
            if normalized in self.visited_urls:
                return False
            if not self.should_process_url(normalized):
                return False
            self.visited_urls.add(normalized)
            self.page_count += 1
            count = self.page_count
            depth = self.url_depths.get(normalized, 0)

        print(f"[{depth}] ({count}/{self.MAX_PAGES}) Crawling: {url}")

        MAX_RETRIES = 3
        response = None
        for attempt in range(MAX_RETRIES):
            try:
                response = self.session.get(url, timeout=30)
                # Fail fast on 404, don't retry
                if response.status_code == 404:
                    response.raise_for_status()
                    
                response.raise_for_status()
                break  # Success
            except requests.RequestException as e:
                # If it's a client error (e.g. 404), don't retry
                if hasattr(e, 'response') and e.response is not None and 400 <= e.response.status_code < 500 and e.response.status_code != 429:
                    with self._lock:
                        self.error_count += 1
                    print(f"  -> Error: {e}")
                    return False
                    
                if attempt < MAX_RETRIES - 1:
                    wait_time = 2 * (attempt + 1)
                    print(f"  -> Retry {attempt+1}/{MAX_RETRIES} for {url} (Wait {wait_time}s) due to error...")
                    time.sleep(wait_time)
                    continue
                else:
                    with self._lock:
                        self.error_count += 1
                    print(f"  -> Error after {MAX_RETRIES} tries: {e}")
                    return False
        # Fix encoding: use apparent_encoding if requests guessed wrong
        if response.encoding and response.encoding.lower() == 'iso-8859-1':
            response.encoding = response.apparent_encoding

        content_type = response.headers.get("Content-Type", "text/html").split(";")[0].strip()
        # Detect extension from URL for fallback
        url_ext = os.path.splitext(urlparse(url).path)[1].lower()

        # Generate filename from URL
        parsed = urlparse(url)
        filename = re.sub(r"[^\w\-]", "_", parsed.path.strip("/").replace("/", "_") or "index")

        content = None
        title = "Untitled"
        raw_html = None
        sub_links = []

        # --- 1. HTML ---
        if "text/html" in content_type and url_ext not in self.BINARY_EXTENSIONS:
            soup = BeautifulSoup(response.text, "html.parser")
            title = self.extract_title(soup)
            raw_html = response.text
            raw_content = self.html_to_markdown(response.text, url)
            content = self.extract_body_content(raw_content)
            # Collect sub-links for BFS queue
            if self.config.recursive and depth < self.config.max_depth:
                sub_links = self.extract_links(soup, url)

        # --- 2. Binary document (PDF/DOCX/XLSX/PPTX) ---
        elif content_type in self.BINARY_CONVERTERS or url_ext in self.BINARY_EXTENSIONS:
            converter_name = None
            if content_type in self.BINARY_CONVERTERS:
                converter_name, _ = self.BINARY_CONVERTERS[content_type]
            else:
                # Fallback: match by URL extension
                for ct, (cn, ext) in self.BINARY_CONVERTERS.items():
                    if ext == url_ext:
                        converter_name = cn
                        break
            if converter_name:
                try:
                    converter = getattr(self, converter_name)
                    content = converter(response.content)
                    # Use filename as title for binary docs
                    title = os.path.basename(parsed.path) or filename
                except Exception as e:
                    print(f"  -> Conversion error ({url_ext}): {e}")
                    return False
            else:
                print(f"  -> Skipped (unsupported binary: {content_type})")
                return False

        # --- 3. Plain text types (txt/csv/json/xml/yaml/md/rst/tex/log) ---
        elif content_type in self.PLAINTEXT_TYPES or url_ext in self.PLAINTEXT_EXTENSIONS:
            text = response.text.strip()
            title = os.path.basename(parsed.path) or filename
            # Wrap structured text in code blocks for readability
            if url_ext in {'.json', '.xml', '.yaml', '.yml', '.csv'}:
                lang = url_ext.lstrip('.')
                if lang == 'yml':
                    lang = 'yaml'
                content = f"```{lang}\n{text}\n```"
            else:
                content = text

        else:
            # Unsupported content type, skip
            print(f"  -> Skipped (unsupported: {content_type})")
            return False

        if content is None:
            return False

        # Prepend source URL
        content = f"> 来源: {url}\n\n{content}"

        # Save content (thread-safe via lock)
        with self._lock:
            status = self.storage.save_content(
                filename=filename,
                content=content,
                source_url=url,
                title=title,
                content_type="text/markdown",
                raw_html=raw_html
            )

            if status == "new":
                full_name = self.current_subdir + '/' + filename
                self.new_files.append(full_name)
                print(f"  -> New: {full_name}")
            elif status == "updated":
                full_name = self.current_subdir + '/' + filename
                self.updated_files.append(full_name)
                print(f"  -> Updated: {full_name}")
            elif status == "unchanged":
                self.unchanged_count += 1
                full_name = self.current_subdir + '/' + filename
                print(f"  -> Unchanged: {full_name}")

        # Add sub-links to BFS queue (only from HTML pages)
        if sub_links:
            with self._lock:
                for link in sub_links:
                    link_norm = self._normalize_url(link)
                    if link_norm not in self.visited_urls and link_norm not in self.url_depths:
                        self.url_depths[link_norm] = depth + 1
                        self.bfs_queue.append(link)

        return True

        # Exception handled inside retry loop

    def _get_delay(self) -> float:
        """Get current delay, checking for real-time updates from .crawl_delay file"""
        try:
            delay_file = os.path.join(os.getcwd(), '.crawl_delay')
            if os.path.exists(delay_file):
                with open(delay_file, 'r') as f:
                    val = float(f.read().strip())
                    if val != self.config.delay:
                        print(f"  [delay updated: {self.config.delay}s -> {val}s]")
                        self.config.delay = val
                    return val
        except (ValueError, IOError):
            pass
        return self.config.delay

    def _crawl_worker(self, url: str):
        """Worker function for concurrent crawling"""
        self.process_page(url)
        time.sleep(self._get_delay())

    def _get_url_subdir(self, url: str) -> str:
        """Extract domain as subdirectory name from URL"""
        parsed = urlparse(url)
        return parsed.netloc or 'default'

    def pre_crawl(self) -> dict:
        """Pre-crawl: BFS discovery of all URLs without downloading/storing content.
        Returns depth statistics for progress estimation."""
        print(f"Pre-crawl: discovering URLs from {len(self.config.urls)} seed(s) ({self.CONCURRENT_WORKERS} workers)")
        self.base_urls = list(self.config.urls)

        discovered = {}  # url -> depth
        visited = set()
        lock = threading.Lock()
        self._pre_crawl_stopped = False  # graceful stop flag

        def _handle_stop(signum, frame):
            print("\n[pre-crawl] Stop requested, saving current stats...")
            self._pre_crawl_stopped = True

        import signal
        signal.signal(signal.SIGTERM, _handle_stop)

        def sniff_worker(url: str, depth: int):
            """Fetch HTML page, extract links. Skip binary files."""
            norm = self._normalize_url(url)
            with lock:
                if norm in visited:
                    return []
                visited.add(norm)

            # Check if URL passes scope + extension filter
            if not self.should_process_url(norm):
                return []

            # For binary file URLs, just count them - don't fetch
            url_ext = os.path.splitext(urlparse(url).path)[1].lower()
            if url_ext in self.BINARY_EXTENSIONS:
                with lock:
                    discovered[norm] = depth
                count = len(discovered)
                print(f"  [pre-crawl] depth={depth} found={count} (doc) {url}")
                return []

            # Fetch HTML pages to extract links
            try:
                resp = self.session.get(url, timeout=15)
                resp.raise_for_status()
            except Exception:
                return []

            content_type = resp.headers.get("Content-Type", "").split(";")[0].strip()
            if "text/html" not in content_type:
                # Non-HTML text file - count it but don't parse for links
                with lock:
                    discovered[norm] = depth
                count = len(discovered)
                print(f"  [pre-crawl] depth={depth} found={count} {url}")
                return []

            # HTML page: count it and extract links
            with lock:
                discovered[norm] = depth
            count = len(discovered)
            print(f"  [pre-crawl] depth={depth} found={count} {url}")

            if resp.encoding and resp.encoding.lower() == 'iso-8859-1':
                resp.encoding = resp.apparent_encoding
            soup = BeautifulSoup(resp.text, "html.parser")
            child_links = []
            for link in self.extract_links(soup, url):
                link_norm = self._normalize_url(link)
                with lock:
                    if link_norm not in visited and link_norm not in discovered:
                        child_links.append((link, depth + 1))
            return child_links

        def _build_result() -> dict:
            """Build statistics from current discovered state."""
            max_depth = max(discovered.values()) if discovered else 0
            urls_per_depth = {}
            for d in range(max_depth + 1):
                urls_per_depth[str(d)] = sum(1 for v in discovered.values() if v == d)
            by_depth = {}
            cumulative = 0
            for d in range(max_depth + 1):
                cumulative += urls_per_depth.get(str(d), 0)
                by_depth[str(d)] = cumulative
            return {
                "total": len(discovered),
                "by_depth": by_depth,
                "urls_per_depth": urls_per_depth,
                "max_depth": max_depth,
                "urls": list(discovered.keys()),
            }

        try:
            for seed_url in self.config.urls:
                if self._pre_crawl_stopped:
                    break
                print(f"\nPre-crawling from: {seed_url}")
                queue = deque([(seed_url, 0)])

                with ThreadPoolExecutor(max_workers=self.CONCURRENT_WORKERS) as executor:
                    active = set()

                    while True:
                        if self._pre_crawl_stopped:
                            executor.shutdown(wait=False, cancel_futures=True)
                            break
                        # Submit from queue
                        while queue and len(active) < self.CONCURRENT_WORKERS:
                            url, depth = queue.popleft()
                            future = executor.submit(sniff_worker, url, depth)
                            future._depth = depth
                            active.add(future)

                        if not active:
                            if not queue:
                                break
                            continue

                        # Wait for one completion
                        done = set()
                        for f in as_completed(active):
                            done.add(f)
                            break
                        active -= done

                        for f in done:
                            try:
                                child_links = f.result()
                                for link, d in child_links:
                                    queue.append((link, d))
                            except Exception:
                                pass

                        time.sleep(0.05)  # light throttle

        except Exception as e:
            print(f"Pre-crawl error: {e}")

        result = _build_result()
        was_stopped = self._pre_crawl_stopped

        if was_stopped:
            print(f"\nPre-crawl stopped early: {len(discovered)} URLs, max depth {result['max_depth']}")
        else:
            print(f"\nPre-crawl complete: {len(discovered)} URLs, max depth {result['max_depth']}")
        for d in range(result['max_depth'] + 1):
            print(f"  depth {d}: {result['urls_per_depth'].get(str(d), 0)} URLs (cumulative: {result['by_depth'].get(str(d), 0)})")

        return result

    def crawl(self) -> CrawlResult:
        """Start crawling from configured URLs using BFS with concurrent workers"""
        print(f"Starting crawl with {len(self.config.urls)} URLs ({self.CONCURRENT_WORKERS} workers)")
        print(f"Output directory: {self.config.output_dir}")
        self.base_urls = list(self.config.urls)  # set scope to configured URLs

        try:
            for url in self.config.urls:
                # Create subdirectory per target URL domain
                subdir = self._get_url_subdir(url)
                self.current_subdir = subdir
                sub_output = os.path.join(self.config.output_dir, subdir)
                self.storage = StorageManager(sub_output, self.config.enable_meta)
                
                print(f"\nCrawling from: {url} -> {subdir}/")
                
                # Initialize BFS with seed URL
                self.bfs_queue = deque([url])
                self.url_depths = {self._normalize_url(url): 0}
                
                # BFS with thread pool
                with ThreadPoolExecutor(max_workers=self.CONCURRENT_WORKERS) as executor:
                    active_futures = set()
                    
                    while True:
                        # Submit new tasks from the queue
                        with self._lock:
                            while self.bfs_queue and len(active_futures) < self.CONCURRENT_WORKERS and self.page_count < self.MAX_PAGES:
                                next_url = self.bfs_queue.popleft()
                                future = executor.submit(self._crawl_worker, next_url)
                                active_futures.add(future)
                        
                        if not active_futures:
                            # No active work and queue is empty - done
                            with self._lock:
                                if not self.bfs_queue:
                                    break
                                else:
                                    continue
                        
                        # Wait for at least one to complete
                        done = set()
                        for f in as_completed(active_futures):
                            done.add(f)
                            break  # Process one at a time to refill queue quickly
                        
                        active_futures -= done
                        
                        # Check if we hit the limit
                        with self._lock:
                            if self.page_count >= self.MAX_PAGES:
                                # Cancel remaining and drain
                                for f in active_futures:
                                    f.cancel()
                                break
                
                # Finalize this URL's storage
                self.storage.finalize(self.new_files, self.updated_files, [])

            return CrawlResult(
                success=True,
                new_files=self.new_files,
                updated_files=self.updated_files,
                deleted_files=[],
                message=f"Crawl completed. Total: {self.page_count}, New: {len(self.new_files)}, Updated: {len(self.updated_files)}, Unchanged: {self.unchanged_count}, Errors: {self.error_count}"
            )

        except Exception as e:
            return CrawlResult(
                success=False,
                new_files=self.new_files,
                updated_files=self.updated_files,
                deleted_files=[],
                message=f"Crawl failed: {str(e)}"
            )


def main():
    """CLI entry point"""
    if len(sys.argv) < 2:
        print("Usage: python crawler.py <config_path> [--pre-crawl]")
        sys.exit(1)

    config_path = sys.argv[1]
    pre_crawl_mode = '--pre-crawl' in sys.argv

    config = CrawlerConfig.from_yaml(config_path)
    crawler = WebCrawler(config)

    if pre_crawl_mode:
        result = crawler.pre_crawl()
        print("\n=== RESULT ===")
        print(json.dumps(result, ensure_ascii=False))
    else:
        result = crawler.crawl()
        print("\n=== RESULT ===")
        print(json.dumps(asdict(result), ensure_ascii=False))


if __name__ == "__main__":
    main()
